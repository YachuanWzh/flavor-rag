"""Shared structured-document IR for non-PDF enterprise formats."""

from __future__ import annotations

import base64
import csv
import io
import json
import mimetypes
import re
import uuid
from dataclasses import dataclass, field
from enum import Enum
from typing import Any

from app.ingestion.pdf.models import PdfAsset, render_markdown_table


class BlockType(str, Enum):
    HEADING = "HEADING"
    PARAGRAPH = "PARAGRAPH"
    LIST = "LIST"
    TABLE = "TABLE"
    CODE = "CODE"
    IMAGE = "IMAGE"


@dataclass
class StructuredBlock:
    block_id: str
    block_type: BlockType
    content: str = ""
    embedding_text: str = ""
    outline_path: list[str] = field(default_factory=list)
    table_headers: list[str] = field(default_factory=list)
    table_rows: list[list[str]] = field(default_factory=list)
    asset_ids: list[str] = field(default_factory=list)
    location: dict[str, Any] = field(default_factory=dict)
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class StructuredDocument:
    source_file: str
    file_type: str
    blocks: list[StructuredBlock] = field(default_factory=list)
    assets: list[PdfAsset] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_markdown(self) -> str:
        parts: list[str] = []
        for block in self.blocks:
            if block.block_type == BlockType.TABLE:
                parts.append(render_markdown_table(block.table_headers, block.table_rows))
            elif block.block_type == BlockType.HEADING:
                parts.append(f"# {block.content}")
            elif block.block_type == BlockType.CODE:
                parts.append(f"```\n{block.content}\n```")
            else:
                parts.append(block.content)
        return "\n\n".join(part for part in parts if part.strip())


def _id(prefix: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex[:16]}"


def _clean(value: Any) -> str:
    return " ".join(str(value or "").replace("\r", "\n").split())


def _markdown_table_cells(line: str) -> list[str]:
    value = line.strip()
    if "|" not in value:
        return []
    if value.startswith("|"):
        value = value[1:]
    if value.endswith("|") and not value.endswith("\\|"):
        value = value[:-1]
    return [
        _clean(cell.replace("\\|", "|"))
        for cell in re.split(r"(?<!\\)\|", value)
    ]


def _is_markdown_table_separator(line: str) -> bool:
    cells = _markdown_table_cells(line)
    return bool(cells) and all(
        re.fullmatch(r":?-{3,}:?", cell.replace(" ", ""))
        for cell in cells
    )


def parse_text_document(text: str, source_file: str, file_type: str) -> StructuredDocument:
    blocks: list[StructuredBlock] = []
    outline: list[str] = []
    in_code = False
    code_lines: list[str] = []
    paragraph: list[str] = []

    def flush_paragraph():
        if not paragraph:
            return
        value = "\n".join(paragraph).strip()
        paragraph.clear()
        if value:
            blocks.append(
                StructuredBlock(
                    _id("paragraph"),
                    BlockType.PARAGRAPH,
                    value,
                    value,
                    list(outline),
                )
            )

    lines = text.splitlines()
    line_index = 0
    while line_index < len(lines):
        line = lines[line_index]
        if line.strip().startswith("```"):
            if in_code:
                blocks.append(
                    StructuredBlock(
                        _id("code"),
                        BlockType.CODE,
                        "\n".join(code_lines).strip(),
                        "\n".join(code_lines).strip(),
                        list(outline),
                    )
                )
                code_lines.clear()
            else:
                flush_paragraph()
            in_code = not in_code
            line_index += 1
            continue
        if in_code:
            code_lines.append(line)
            line_index += 1
            continue
        if (
            line_index + 1 < len(lines)
            and _markdown_table_cells(line)
            and _is_markdown_table_separator(lines[line_index + 1])
        ):
            flush_paragraph()
            headers = _markdown_table_cells(line)
            table_start = line_index
            line_index += 2
            rows: list[list[str]] = []
            while line_index < len(lines) and lines[line_index].strip():
                cells = _markdown_table_cells(lines[line_index])
                if not cells:
                    break
                rows.append(
                    (cells + [""] * len(headers))[: len(headers)]
                )
                line_index += 1
            blocks.append(
                StructuredBlock(
                    _id("table"),
                    BlockType.TABLE,
                    outline_path=list(outline),
                    table_headers=headers,
                    table_rows=rows,
                    location={
                        "lineStart": table_start + 1,
                        "lineEnd": line_index,
                    },
                    metadata={"format": "markdown"},
                )
            )
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
        if heading:
            flush_paragraph()
            level = len(heading.group(1))
            outline[:] = outline[: level - 1]
            outline.append(heading.group(2).strip())
            blocks.append(
                StructuredBlock(
                    _id("heading"),
                    BlockType.HEADING,
                    outline[-1],
                    outline[-1],
                    list(outline),
                    metadata={"level": level},
                )
            )
            line_index += 1
            continue
        if re.match(r"^\s*(?:[-*+]|\d+[.)])\s+", line):
            flush_paragraph()
            value = re.sub(r"^\s*(?:[-*+]|\d+[.)])\s+", "", line).strip()
            blocks.append(
                StructuredBlock(
                    _id("list"),
                    BlockType.LIST,
                    f"- {value}",
                    value,
                    list(outline),
                )
            )
            line_index += 1
            continue
        if not line.strip():
            flush_paragraph()
        else:
            paragraph.append(line)
        line_index += 1
    flush_paragraph()
    if code_lines:
        blocks.append(
            StructuredBlock(
                _id("code"),
                BlockType.CODE,
                "\n".join(code_lines),
                "\n".join(code_lines),
                list(outline),
            )
        )
    return StructuredDocument(source_file, file_type, blocks)


async def parse_clipboard_document(
    content: bytes,
    source_file: str,
) -> StructuredDocument:
    """Parse the internal text-and-image clipboard bundle format."""
    payload = json.loads(content.decode("utf-8"))
    if payload.get("version") != 1:
        raise ValueError("Unsupported clipboard document version")

    text = str(payload.get("content") or "")
    image_payloads = payload.get("images") or []
    if not isinstance(image_payloads, list):
        raise ValueError("Clipboard document images must be a list")

    from app.ingestion.pdf.vlm import get_image_describer

    describer = get_image_describer()
    image_info_by_id: dict[str, dict[str, Any]] = {}
    assets: list[PdfAsset] = []
    for index, image in enumerate(image_payloads, start=1):
        image_id = str(image.get("id") or "")
        if not image_id:
            raise ValueError(f"Clipboard image at index {index} has no id")
        try:
            image_bytes = base64.b64decode(image["data"], validate=True)
        except Exception as exc:
            raise ValueError(f"Invalid clipboard image at index {index}") from exc
        mime_type = str(image.get("mimeType") or "application/octet-stream")
        filename = str(image.get("filename") or f"clipboard-image-{index}")
        alt = str(image.get("alt") or filename).strip()
        try:
            description = await describer.describe(
                image_bytes,
                mime_type,
                context=f"剪贴板文档={source_file}; 图片序号={index}",
            )
            description_status = "enriched" if description else "missing"
        except Exception:
            description = ""
            description_status = "failed"
        asset = PdfAsset.from_bytes(
            page_no=index,
            bbox=None,
            filename=filename,
            mime_type=mime_type,
            data=image_bytes,
            description=description,
            metadata={
                "source": "clipboard",
                "description_status": description_status,
            },
        )
        assets.append(asset)
        image_info_by_id[image_id] = {
            "index": index,
            "asset": asset,
            "alt": alt,
            "description": description,
            "description_status": description_status,
            "mime_type": mime_type,
        }

    marker_pattern = re.compile(
        r"!\[([^\]]*)\]\(clipboard-image://([A-Za-z0-9._-]+)\)"
    )

    def context_text(value: str, *, tail: bool = False) -> str:
        value = marker_pattern.sub("", value)
        value = " ".join(value.split())
        return value[-500:] if tail else value[:500]

    blocks: list[StructuredBlock] = []
    current_outline: list[str] = []

    def append_text_segment(segment: str) -> None:
        nonlocal current_outline
        parsed = parse_text_document(segment, source_file, "clipboard")
        for block in parsed.blocks:
            if block.block_type == BlockType.HEADING:
                level = max(1, int(block.metadata.get("level") or 1))
                current_outline = current_outline[: level - 1]
                current_outline.append(block.content)
                block.outline_path = list(current_outline)
            elif not block.outline_path:
                block.outline_path = list(current_outline)
            blocks.append(block)

    def append_image_block(
        image_id: str,
        marker_alt: str,
        *,
        offset: int,
        before: str,
        after: str,
    ) -> bool:
        info = image_info_by_id.get(image_id)
        if info is None:
            return False
        asset = info["asset"]
        alt = marker_alt.strip() or info["alt"]
        visible = info["description"] or alt or asset.filename
        embedding_parts = []
        if current_outline:
            embedding_parts.append(f"章节: {' > '.join(current_outline)}")
        if before:
            embedding_parts.append(f"图片前文: {before}")
        embedding_parts.append(f"图片内容: {visible}")
        if after:
            embedding_parts.append(f"图片后文: {after}")
        blocks.append(
            StructuredBlock(
                _id("image"),
                BlockType.IMAGE,
                f"{info['description']}\n\n![{alt}](asset://{asset.asset_id})".strip(),
                "\n".join(embedding_parts),
                list(current_outline),
                asset_ids=[asset.asset_id],
                location={
                    "clipboardImage": info["index"],
                    "charOffset": offset,
                },
                metadata={
                    "source": "clipboard",
                    "descriptionStatus": info["description_status"],
                    "mimeType": info["mime_type"],
                    "contextBefore": before,
                    "contextAfter": after,
                },
            )
        )
        return True

    referenced_ids: set[str] = set()
    cursor = 0
    for marker in marker_pattern.finditer(text):
        append_text_segment(text[cursor:marker.start()])
        image_id = marker.group(2)
        if append_image_block(
            image_id,
            marker.group(1),
            offset=marker.start(),
            before=context_text(text[:marker.start()], tail=True),
            after=context_text(text[marker.end():]),
        ):
            referenced_ids.add(image_id)
        cursor = marker.end()
    append_text_segment(text[cursor:])

    # Preserve assets even when a source application supplied image bytes but
    # omitted the corresponding HTML marker.
    for image_id, info in image_info_by_id.items():
        if image_id not in referenced_ids:
            append_image_block(
                image_id,
                info["alt"],
                offset=len(text),
                before=context_text(text, tail=True),
                after="",
            )

    document = StructuredDocument(
        source_file,
        "clipboard",
        blocks,
        assets,
        metadata={
            "imageCount": len(assets),
            "sourceFormat": "clipboard",
        },
    )
    return document


def parse_csv_document(content: bytes, source_file: str) -> StructuredDocument:
    text = content.decode("utf-8-sig", errors="replace")
    rows = [list(row) for row in csv.reader(io.StringIO(text))]
    headers = [_clean(item) for item in (rows[0] if rows else [])]
    data = [[_clean(item) for item in row] for row in rows[1:]]
    block = StructuredBlock(
        _id("table"),
        BlockType.TABLE,
        table_headers=headers,
        table_rows=data,
        location={"rowStart": 2, "rowEnd": len(rows)},
        metadata={"sheet": source_file, "format": "csv"},
    )
    return StructuredDocument(source_file, "csv", [block])


def parse_docx_document(content: bytes, source_file: str) -> StructuredDocument:
    from docx import Document

    doc = Document(io.BytesIO(content))
    blocks: list[StructuredBlock] = []
    outline: list[str] = []
    paragraph_index = 0
    for paragraph in doc.paragraphs:
        value = paragraph.text.strip()
        if not value:
            continue
        paragraph_index += 1
        style = (paragraph.style.name if paragraph.style else "").lower()
        heading = re.search(r"heading\s*(\d+)|标题\s*(\d+)", style)
        if heading:
            level = int(next(group for group in heading.groups() if group) or 1)
            outline[:] = outline[: level - 1]
            outline.append(value)
            block_type = BlockType.HEADING
        elif "list" in style or "列表" in style:
            block_type = BlockType.LIST
        else:
            block_type = BlockType.PARAGRAPH
        blocks.append(
            StructuredBlock(
                _id(block_type.value.lower()),
                block_type,
                f"- {value}" if block_type == BlockType.LIST else value,
                value,
                list(outline),
                location={"paragraph": paragraph_index},
                metadata={"style": paragraph.style.name if paragraph.style else ""},
            )
        )
    for table_index, table in enumerate(doc.tables, start=1):
        rows = [[_clean(cell.text) for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        blocks.append(
            StructuredBlock(
                _id("table"),
                BlockType.TABLE,
                outline_path=list(outline),
                table_headers=rows[0],
                table_rows=rows[1:],
                location={"table": table_index},
                metadata={"format": "docx"},
            )
        )
    return StructuredDocument(
        source_file,
        "docx",
        blocks,
        metadata={"paragraphCount": len(doc.paragraphs), "tableCount": len(doc.tables)},
    )


def parse_xlsx_document(content: bytes, source_file: str) -> StructuredDocument:
    import openpyxl

    values_wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    formula_wb = openpyxl.load_workbook(io.BytesIO(content), data_only=False)
    blocks: list[StructuredBlock] = []
    for sheet_name in values_wb.sheetnames:
        ws = values_wb[sheet_name]
        formula_ws = formula_wb[sheet_name]
        blocks.append(
            StructuredBlock(
                _id("heading"),
                BlockType.HEADING,
                sheet_name,
                sheet_name,
                [sheet_name],
                location={"sheet": sheet_name},
                metadata={"level": 1},
            )
        )
        rows = [
            [_clean(cell.value) for cell in row]
            for row in ws.iter_rows()
            if any(cell.value is not None for cell in row)
        ]
        if not rows:
            continue
        hyperlinks: dict[str, str] = {}
        formulas: dict[str, str] = {}
        for row in ws.iter_rows():
            for cell in row:
                if cell.hyperlink and cell.hyperlink.target:
                    hyperlinks[cell.coordinate] = cell.hyperlink.target
                raw = formula_ws[cell.coordinate].value
                if isinstance(raw, str) and raw.startswith("="):
                    formulas[cell.coordinate] = raw
        blocks.append(
            StructuredBlock(
                _id("table"),
                BlockType.TABLE,
                outline_path=[sheet_name],
                table_headers=rows[0],
                table_rows=rows[1:],
                location={
                    "sheet": sheet_name,
                    "rowStart": 1,
                    "rowEnd": ws.max_row,
                },
                metadata={
                    "format": "xlsx",
                    "mergedCells": [str(item) for item in ws.merged_cells.ranges],
                    "hyperlinks": hyperlinks,
                    "formulas": formulas,
                },
            )
        )
    return StructuredDocument(
        source_file,
        "xlsx",
        blocks,
        metadata={"sheetCount": len(values_wb.sheetnames)},
    )


def parse_pptx_document(content: bytes, source_file: str) -> StructuredDocument:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    blocks: list[StructuredBlock] = []
    for slide_no, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else f"第 {slide_no} 页"
        blocks.append(
            StructuredBlock(
                _id("heading"),
                BlockType.HEADING,
                title,
                title,
                [title],
                location={"slide": slide_no},
                metadata={"level": 1},
            )
        )
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape == slide.shapes.title:
                continue
            if getattr(shape, "has_table", False):
                rows = [
                    [_clean(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    blocks.append(
                        StructuredBlock(
                            _id("table"),
                            BlockType.TABLE,
                            outline_path=[title],
                            table_headers=rows[0],
                            table_rows=rows[1:],
                            location={"slide": slide_no, "shape": shape_index},
                            metadata={"format": "pptx"},
                        )
                    )
            elif getattr(shape, "has_text_frame", False):
                value = shape.text_frame.text.strip()
                if value:
                    blocks.append(
                        StructuredBlock(
                            _id("paragraph"),
                            BlockType.PARAGRAPH,
                            value,
                            value,
                            [title],
                            location={"slide": slide_no, "shape": shape_index},
                        )
                    )
    return StructuredDocument(
        source_file,
        "pptx",
        blocks,
        metadata={"slideCount": len(presentation.slides)},
    )


def parse_html_document(content: bytes, source_file: str) -> StructuredDocument:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    blocks: list[StructuredBlock] = []
    outline: list[str] = []
    for element in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "pre", "table", "img"]):
        if element.find_parent(["table", "pre"]) and element.name not in {"table", "pre"}:
            continue
        value = element.get_text(" ", strip=True)
        if element.name.startswith("h"):
            level = int(element.name[1])
            outline[:] = outline[: level - 1]
            outline.append(value)
            blocks.append(
                StructuredBlock(
                    _id("heading"), BlockType.HEADING, value, value,
                    list(outline), metadata={"level": level},
                )
            )
        elif element.name == "table":
            rows = [
                [_clean(cell.get_text(" ", strip=True)) for cell in row.find_all(["th", "td"])]
                for row in element.find_all("tr")
            ]
            if rows:
                blocks.append(
                    StructuredBlock(
                        _id("table"), BlockType.TABLE,
                        outline_path=list(outline),
                        table_headers=rows[0], table_rows=rows[1:],
                        metadata={"format": "html"},
                    )
                )
        elif element.name == "img":
            alt = element.get("alt") or element.get("src") or "网页图片"
            blocks.append(
                StructuredBlock(
                    _id("image"), BlockType.IMAGE, alt, alt,
                    list(outline), metadata={"src": element.get("src", "")},
                )
            )
        elif value:
            block_type = {
                "li": BlockType.LIST,
                "pre": BlockType.CODE,
            }.get(element.name, BlockType.PARAGRAPH)
            blocks.append(
                StructuredBlock(
                    _id(block_type.value.lower()),
                    block_type,
                    f"- {value}" if block_type == BlockType.LIST else value,
                    value,
                    list(outline),
                )
            )
    return StructuredDocument(source_file, "html", blocks)


async def parse_image_document(
    content: bytes,
    source_file: str,
    file_type: str,
) -> StructuredDocument:
    from app.ingestion.pdf.vlm import get_image_describer

    mime_type = mimetypes.guess_type(source_file)[0] or f"image/{file_type}"
    description = await get_image_describer().describe(
        content,
        mime_type,
        context=source_file,
    )
    asset = PdfAsset.from_bytes(
        page_no=1,
        bbox=None,
        filename=source_file,
        mime_type=mime_type,
        data=content,
        description=description,
    )
    visible = description or source_file
    block = StructuredBlock(
        _id("image"),
        BlockType.IMAGE,
        f"{description}\n\n![{source_file}](asset://{asset.asset_id})".strip(),
        visible,
        asset_ids=[asset.asset_id],
        location={"image": 1},
        metadata={
            "descriptionStatus": "enriched" if description else "missing",
            "mimeType": mime_type,
        },
    )
    return StructuredDocument(source_file, file_type, [block], [asset])


class GenericStructuredChunker:
    """Block-aware chunker retaining format-specific provenance."""

    def __init__(self, *, target_chars: int = 800, table_max_rows: int = 20):
        self.target_chars = max(128, target_chars)
        self.table_max_rows = max(1, table_max_rows)

    def chunk(self, document: StructuredDocument) -> list[dict]:
        chunks: list[dict] = []
        outline: list[str] = []
        text_buffer: list[StructuredBlock] = []

        def flush_text():
            nonlocal text_buffer
            if not text_buffer:
                return
            content = "\n\n".join(block.content for block in text_buffer if block.content)
            embedding = "\n".join(
                [f"章节: {' > '.join(outline)}"] if outline else []
            ) + ("\n" if outline else "") + content
            chunks.append(
                self._record(
                    chunks,
                    content,
                    embedding,
                    "PARAGRAPH",
                    {
                        "sourceBlockIds": [block.block_id for block in text_buffer],
                        "outlinePath": list(outline),
                        "locations": [block.location for block in text_buffer],
                        "sourceFormat": document.file_type,
                    },
                )
            )
            text_buffer = []

        for block in document.blocks:
            if block.block_type == BlockType.HEADING:
                flush_text()
                outline = list(block.outline_path or [block.content])
                continue
            if block.block_type == BlockType.TABLE:
                flush_text()
                for start in range(0, len(block.table_rows) or 1, self.table_max_rows):
                    rows = block.table_rows[start : start + self.table_max_rows]
                    content = render_markdown_table(block.table_headers, rows)
                    section = " > ".join(block.outline_path or outline)
                    embedding_lines = [f"章节: {section}"] if section else []
                    for row in rows:
                        pairs = [
                            f"{header}: {row[index] if index < len(row) else ''}"
                            for index, header in enumerate(block.table_headers)
                        ]
                        embedding_lines.append("；".join(pairs))
                    chunks.append(
                        self._record(
                            chunks,
                            content,
                            "\n".join(embedding_lines),
                            "TABLE",
                            {
                                **block.metadata,
                                "sourceBlockIds": [block.block_id],
                                "outlinePath": block.outline_path or outline,
                                "location": block.location,
                                "rowStart": start + 1,
                                "rowEnd": start + len(rows),
                                "sourceFormat": document.file_type,
                            },
                        )
                    )
                continue
            if block.block_type in {BlockType.IMAGE, BlockType.CODE}:
                flush_text()
                chunks.append(
                    self._record(
                        chunks,
                        block.content,
                        block.embedding_text or block.content,
                        block.block_type.value,
                        {
                            **block.metadata,
                            "sourceBlockIds": [block.block_id],
                            "outlinePath": block.outline_path or outline,
                            "location": block.location,
                            "sourceFormat": document.file_type,
                        },
                        asset_ids=block.asset_ids,
                    )
                )
                continue
            prospective = sum(len(item.content) for item in text_buffer) + len(block.content)
            if text_buffer and prospective > self.target_chars:
                flush_text()
            text_buffer.append(block)
        flush_text()
        return chunks

    @staticmethod
    def _record(
        chunks: list[dict],
        content: str,
        embedding: str,
        block_type: str,
        metadata: dict,
        *,
        asset_ids: list[str] | None = None,
    ) -> dict:
        return {
            "chunk_index": len(chunks),
            "content": content,
            "embedding_content": embedding,
            "char_count": len(content),
            "block_type": block_type,
            "page_start": None,
            "page_end": None,
            "bbox_json": [],
            "metadata_json": metadata,
            "asset_ids": list(asset_ids or []),
        }
