"""Parser node — parses raw document bytes into plain text."""

from __future__ import annotations

import os
import tempfile

from app.config.logging_config import get_logger
from app.ingestion.nodes.base import IngestionContext, NodeResult

_log = get_logger("flavorag.ingestion.parser")


class ParserNode:
    """Parse raw content bytes into plain text using Tika-like fallback chain.

    Supports: PDF, DOCX, XLSX, PPTX, Markdown, TXT, HTML, CSV.

    Settings:
        parser_type (str): "tika" (preferred) or "native" (default).
        tika_url (str): Tika server URL for tika mode.
    """

    NODE_TYPE = "parser"

    def __init__(self, structured_pdf_parser=None):
        if structured_pdf_parser is None:
            from app.ingestion.pdf.parser import StructuredPdfParser
            structured_pdf_parser = StructuredPdfParser()
        self.structured_pdf_parser = structured_pdf_parser

    async def __call__(self, ctx: IngestionContext) -> NodeResult:
        import time
        t0 = time.time()

        try:
            if ctx.raw_content is None:
                raise ValueError("No raw content to parse")

            parser_type = ctx.settings.get("parser_type", "native")

            ext = self._guess_extension(ctx.source_file_name)
            structured = parser_type != "tika"
            if structured:
                if ext == "pdf":
                    document = await self.structured_pdf_parser.parse_bytes(
                        ctx.raw_content,
                        source_file=ctx.source_file_name,
                        document_id=ctx.doc_id,
                    )
                else:
                    document = await self._parse_structured(ctx, ext)
                ctx.parsed_document = document
                ctx.assets = list(document.assets)
                ctx.metadata["structured_document"] = {
                    "format": ext,
                    **document.metadata,
                }
                text = document.to_markdown()
            elif parser_type == "tika":
                text = await self._parse_tika(ctx)
            else:
                text = self._parse_native(ctx)

            ctx.parsed_text = text
            duration_ms = int((time.time() - t0) * 1000)
            _log.info(
                "parser_done",
                doc_id=ctx.doc_id,
                text_len=len(text),
                parser=parser_type,
                took_ms=duration_ms,
            )
            return NodeResult(
                node_type=self.NODE_TYPE,
                status="success",
                duration_ms=duration_ms,
                output={
                    "text_length": len(text),
                    "structured": structured,
                    "page_count": getattr(ctx.parsed_document, "page_count", None) if structured else None,
                    "table_count": sum(
                        getattr(block.block_type, "value", block.block_type) == "TABLE"
                        for block in ctx.parsed_document.blocks
                    ) if structured else 0,
                    "image_count": len(ctx.assets) if structured else 0,
                },
            )
        except Exception as exc:
            duration_ms = int((time.time() - t0) * 1000)
            _log.error("parser_failed", doc_id=ctx.doc_id, error=str(exc))
            return NodeResult(
                node_type=self.NODE_TYPE,
                status="error",
                error_message=str(exc),
                duration_ms=duration_ms,
            )

    async def _parse_tika(self, ctx: IngestionContext) -> str:
        """Send document to Apache Tika server for text extraction."""
        import httpx
        tika_url = ctx.settings.get("tika_url", "http://localhost:9998/tika")
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.put(
                tika_url,
                content=ctx.raw_content,
                headers={"Accept": "text/plain"},
            )
            resp.raise_for_status()
            return resp.text

    async def _parse_structured(self, ctx: IngestionContext, ext: str):
        from app.ingestion.structured import (
            parse_csv_document,
            parse_docx_document,
            parse_html_document,
            parse_image_document,
            parse_pptx_document,
            parse_text_document,
            parse_xlsx_document,
        )

        content = ctx.raw_content or b""
        name = ctx.source_file_name or f"document.{ext}"
        if ext in ("txt", "md", "markdown", "log"):
            return parse_text_document(self._decode_text(content), name, ext)
        if ext == "csv":
            return parse_csv_document(content, name)
        if ext in ("docx", "doc"):
            return parse_docx_document(content, name)
        if ext in ("xlsx", "xls"):
            return parse_xlsx_document(content, name)
        if ext in ("pptx", "ppt"):
            return parse_pptx_document(content, name)
        if ext in ("html", "htm"):
            return parse_html_document(content, name)
        if ext in ("png", "jpg", "jpeg", "webp", "gif", "bmp"):
            return await parse_image_document(content, name, ext)
        return parse_text_document(self._decode_text(content), name, ext or "txt")

    def _parse_native(self, ctx: IngestionContext) -> str:
        """Native Python parsing fallback chain."""
        ext = self._guess_extension(ctx.source_file_name)
        content = ctx.raw_content

        if ext in ("txt", "md", "markdown", "log", "csv"):
            return self._decode_text(content)

        if ext == "pdf":
            return self._parse_pdf(content)

        if ext in ("docx", "doc"):
            return self._parse_docx(content)

        if ext in ("xlsx", "xls"):
            return self._parse_xlsx(content)

        if ext in ("pptx", "ppt"):
            return self._parse_pptx(content)

        if ext in ("html", "htm"):
            return self._parse_html(content)

        # Fallback: try as plain text
        return self._decode_text(content)

    def _decode_text(self, content: bytes) -> str:
        for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace")

    def _parse_pdf(self, content: bytes) -> str:
        try:
            import pymupdf
            tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
            try:
                tmp.write(content)
                tmp.close()
                doc = pymupdf.open(tmp.name)
                pages = [page.get_text() for page in doc]
                return "\n\n".join(pages)
            finally:
                os.unlink(tmp.name)
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                from io import BytesIO
                reader = PdfReader(BytesIO(content))
                pages = [page.extract_text() or "" for page in reader.pages]
                return "\n\n".join(pages)
            except ImportError:
                return self._decode_text(content)

    def _parse_docx(self, content: bytes) -> str:
        try:
            from docx import Document
            from io import BytesIO
            doc = Document(BytesIO(content))
            return "\n\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            return self._decode_text(content)

    def _parse_xlsx(self, content: bytes) -> str:
        try:
            import openpyxl
            from io import BytesIO
            wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    parts.append(" | ".join(str(c) if c is not None else "" for c in row))
            return "\n".join(parts)
        except ImportError:
            return self._decode_text(content)

    def _parse_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"[Slide {i + 1}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n\n".join(parts)
        except ImportError:
            return self._decode_text(content)

    def _parse_html(self, content: bytes) -> str:
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            return self._decode_text(content)

    @staticmethod
    def _guess_extension(filename: str) -> str:
        if not filename:
            return "txt"
        return filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"
